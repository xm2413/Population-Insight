from __future__ import annotations

import sys
from datetime import date
from pathlib import Path

ROOT = Path('/Users/mjy/ClassWork/大三下/Python/Class Project')
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

SRC = ROOT / 'ppt' / '项目进展汇报PPT.pptx'
OUT = ROOT / 'ppt' / '人口趋势与分析管理系统_项目进展汇报.pptx'
CHART_DIR = ROOT / 'static' / 'charts'

prs = Presentation(SRC)

try:
    from app import app
    from population_trend.database import init_database
    from population_trend.visualization_service import generate_charts
    init_database()
    with app.app_context():
        generate_charts()
except Exception as exc:
    print(f'chart generation skipped: {exc}')

ACCENT = RGBColor(31, 94, 168)
DARK = RGBColor(34, 45, 60)
MUTED = RGBColor(88, 96, 110)
WHITE = RGBColor(255, 255, 255)


def iter_all_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:
            yield from iter_all_shapes(shape.shapes)


def iter_text_shapes(slide):
    for shape in iter_all_shapes(slide.shapes):
        if getattr(shape, 'has_text_frame', False) and shape.has_text_frame:
            yield shape


def clean_text(shape):
    return shape.text.strip().replace('\n', '').replace(' ', '')


def set_text(shape, text, size=None, color=None, bold=None, align=None):
    shape.text = text
    shape.text_frame.word_wrap = True
    shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    for p in shape.text_frame.paragraphs:
        if align is not None:
            p.alignment = align
        for run in p.runs:
            if size is not None:
                run.font.size = Pt(size)
            if color is not None:
                run.font.color.rgb = color
            if bold is not None:
                run.font.bold = bold


def set_bullets(shape, title, bullets, title_size=20, body_size=15):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    p.runs[0].font.size = Pt(title_size)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = DARK
    for bullet in bullets:
        para = tf.add_paragraph()
        para.text = bullet
        para.level = 0
        para.space_before = Pt(5)
        para.runs[0].font.size = Pt(body_size)
        para.runs[0].font.color.rgb = MUTED


def replace_exact(slide, mapping):
    for shape in iter_text_shapes(slide):
        raw = clean_text(shape)
        if raw in mapping:
            text, size, color, bold, align = mapping[raw]
            set_text(shape, text, size=size, color=color, bold=bold, align=align)


def add_note(slide, text, left, top, width, height, size=12):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box, text, size=size, color=MUTED, bold=False, align=PP_ALIGN.LEFT)
    box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    return box

# Slide 1
s = prs.slides[0]
for shape in iter_text_shapes(s):
    raw = clean_text(shape)
    if raw == '项目名称':
        set_text(shape, '人口趋势与分析管理系统', size=34, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    elif '数据库与程序设计' in raw:
        set_text(shape, 'Python + Flask + SQLite + Matplotlib', size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    elif any(key in raw for key in ['汇报人', '成员一', '成员二', '成员三', '成员四']):
        set_text(shape, '汇报人：项目小组\n组员：请填入成员姓名', size=15, color=DARK, bold=False, align=PP_ALIGN.LEFT)
    elif '指导老师' in raw:
        set_text(shape, '指导老师：请填入老师姓名', size=15, color=DARK, bold=False, align=PP_ALIGN.LEFT)
    elif '专业' in raw:
        set_text(shape, '课程：数据库与程序设计', size=15, color=DARK, bold=False, align=PP_ALIGN.LEFT)

# Slide 2
replace_exact(prs.slides[1], {
    '项目概括与目标': ('项目概括与目标', 25, DARK, True, PP_ALIGN.LEFT),
    '项目进度完成情况': ('项目进度完成情况', 25, DARK, True, PP_ALIGN.LEFT),
    '下一步的工作计划': ('下一步工作计划', 25, DARK, True, PP_ALIGN.LEFT),
    '项目总结与体会': ('项目总结与体会', 25, DARK, True, PP_ALIGN.LEFT),
})

# Slide 3
replace_exact(prs.slides[2], {
    '项目概括与目标': ('项目概括与目标', 33, WHITE, True, PP_ALIGN.CENTER),
    'PersonalIntroduction': ('Project Overview & Objectives', 17, WHITE, False, PP_ALIGN.CENTER),
})

# Slide 4
s = prs.slides[3]
replace_exact(s, {
    '主要研究方向': ('项目定位', 18, DARK, True, PP_ALIGN.LEFT),
    '硕博阶段成就': ('建设目标', 18, DARK, True, PP_ALIGN.LEFT),
})
for shape in iter_text_shapes(s):
    raw = clean_text(shape)
    if raw.startswith('场景感知'):
        set_bullets(shape, '小型数据分析类信息管理系统', [
            '面向人口数据的录入、维护、筛选、统计和展示',
            '使用 SQLite 完成持久化存储，覆盖增删改查',
            '使用 Matplotlib 输出折线图、柱状图、饼图',
            '新增全球人口统计模块，支持国家与洲维度分析',
        ], title_size=18, body_size=13)
    elif raw.startswith('1.硕博'):
        set_bullets(shape, '课程要求对齐', [
            '独立功能模块超过 8 个：登录、CRUD、查询、导入、导出、统计、可视化、预测、预警、日志等',
            '采用 population_trend 包进行分文件封装',
            '数据库自动初始化并写入地区与全球样例数据',
            '可通过 Web 页面完成现场演示',
        ], title_size=18, body_size=13)

# Slide 5
replace_exact(prs.slides[4], {
    '项目进度完成情况': ('项目进度完成情况', 33, WHITE, True, PP_ALIGN.CENTER),
    'Trialteachingcontent': ('Completed Work & Current Status', 17, WHITE, False, PP_ALIGN.CENTER),
})

# Slide 6
s = prs.slides[5]
replace_exact(s, {
    '主要研究方向': ('已完成功能', 18, DARK, True, PP_ALIGN.LEFT),
    '硕博阶段成就': ('技术完成度', 18, DARK, True, PP_ALIGN.LEFT),
})
for shape in iter_text_shapes(s):
    raw = clean_text(shape)
    if raw.startswith('场景感知'):
        set_bullets(shape, '系统核心模块已经打通', [
            '登录权限：管理员与查看者角色区分',
            '人口数据：新增、查询、修改、删除、筛选、导出',
            '统计分析：地区排名、人口变化、增长率分析',
            '全球统计：国家人口排行、洲际汇总、全球趋势',
            '风险预警：老龄化、负增长、低出生率识别',
        ], title_size=18, body_size=13)
    elif raw.startswith('1.硕博'):
        set_bullets(shape, '代码与数据已具备演示条件', [
            'app.py 负责 Web 路由，业务逻辑拆分到 population_trend 包',
            'database.py 自动建表并初始化样例数据',
            'analysis_service.py 封装统计、预测、预警和全球统计',
            'visualization_service.py 生成 Matplotlib 图表',
            'README 与项目说明书已同步更新',
        ], title_size=18, body_size=13)

# Slide 7
replace_exact(prs.slides[6], {
    '下一步的工作计划': ('下一步工作计划', 33, WHITE, True, PP_ALIGN.CENTER),
    'ResearchContents': ('Next Sprint Plan', 17, WHITE, False, PP_ALIGN.CENTER),
})

# Slide 8
s = prs.slides[7]
replace_exact(s, {'xxxxxxxxxxxx': ('下一阶段任务拆解', 30, DARK, True, PP_ALIGN.CENTER)})
shorts = ['数据完善', '功能联调', '页面优化', '图表检查', '报告整理', '演示排练']
longs = [
    '补充真实或更完整的全球国家 CSV 数据',
    '检查新增、修改、导入、导出流程',
    '统一按钮、表格和移动端布局',
    '确保折线图、柱状图、饼图可正常生成',
    '补充 ER 图、流程图和测试截图',
    '准备 3-5 分钟现场演示路线',
]
short_idx = 0
long_idx = 0
for shape in iter_text_shapes(s):
    raw = clean_text(shape)
    if raw == 'xxx' and shape.height < 500000 and short_idx < len(shorts):
        set_text(shape, shorts[short_idx], size=15, color=DARK, bold=True, align=PP_ALIGN.CENTER)
        short_idx += 1
    elif raw in {'xxxxx', '.........', 'xxxx', 'xxxx。'} and shape.height >= 500000 and long_idx < len(longs):
        set_text(shape, longs[long_idx], size=11, color=MUTED, bold=False, align=PP_ALIGN.CENTER)
        long_idx += 1
# Any leftover short placeholders get filled in order.
for shape in iter_text_shapes(s):
    if clean_text(shape) == 'xxx' and short_idx < len(shorts):
        set_text(shape, shorts[short_idx], size=15, color=DARK, bold=True, align=PP_ALIGN.CENTER)
        short_idx += 1

# Slide 9
replace_exact(prs.slides[8], {
    '项目总结与体会': ('项目总结与体会', 33, WHITE, True, PP_ALIGN.CENTER),
    'FuturePlans': ('Summary & Reflection', 17, WHITE, False, PP_ALIGN.CENTER),
})

# Slide 10
s = prs.slides[9]
# Specific shape names from template avoid duplicate placeholder ambiguity.
specific = {
    '标题 4': ('系统架构与数据流', 30, DARK, True, PP_ALIGN.CENTER),
    '燕尾形 14': ('数据层', 15, WHITE, True, PP_ALIGN.CENTER),
    '燕尾形 17': ('服务层', 15, WHITE, True, PP_ALIGN.CENTER),
    '燕尾形 20': ('分析层', 15, WHITE, True, PP_ALIGN.CENTER),
    '燕尾形 23': ('展示层', 15, WHITE, True, PP_ALIGN.CENTER),
    '矩形 44': ('SQLite 建表\n样例数据\n持久化存储', 13, DARK, False, PP_ALIGN.CENTER),
    '矩形 30': ('record_service\nauth\nvalidators', 13, DARK, False, PP_ALIGN.CENTER),
    '矩形 48': ('统计分析\n线性预测\n风险预警', 13, DARK, False, PP_ALIGN.CENTER),
    '矩形 43': ('Flask 页面\nCSV 导入导出\nMatplotlib 图表', 13, DARK, False, PP_ALIGN.CENTER),
}
for shape in iter_text_shapes(s):
    if shape.name in specific:
        set_text(shape, *specific[shape.name])

# Slide 11
s = prs.slides[10]
for shape in iter_text_shapes(s):
    raw = clean_text(shape)
    if raw.startswith('单击此处'):
        set_bullets(shape, '现场演示与验收要点', [
            '登录：admin / admin123，查看者 viewer / viewer123',
            '人口数据管理：筛选、新增、编辑、删除、CSV 导出',
            '全球统计：按洲筛选，查看国家人口排行与洲际汇总',
            '可视化页面：Matplotlib 输出折线图、柱状图、饼图',
            '项目说明：代码分模块，数据库可持久化，功能数量满足要求',
        ], title_size=23, body_size=16)
chart1 = CHART_DIR / 'global_population_trend.png'
chart2 = CHART_DIR / 'population_ranking.png'
if chart1.exists():
    s.shapes.add_picture(str(chart1), Inches(5.85), Inches(2.25), width=Inches(3.0))
if chart2.exists():
    s.shapes.add_picture(str(chart2), Inches(8.9), Inches(2.25), width=Inches(3.0))
add_note(s, '图表由系统运行时自动生成，可作为可视化模块验收依据。', Inches(5.85), Inches(5.25), Inches(5.8), Inches(0.45), size=11)

# Slide 12
s = prs.slides[11]
for shape in iter_text_shapes(s):
    raw = clean_text(shape)
    if raw == '汇报人：xxx':
        set_text(shape, '汇报人：项目小组', size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    elif raw == '谢谢聆听欢迎同学老师指导':
        set_text(shape, '谢谢聆听\n欢迎老师和同学批评指导', size=34, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    elif raw.startswith('时间：'):
        set_text(shape, f'时间：{date.today().year}年{date.today().month:02d}月{date.today().day:02d}日', size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Fill empty slide number placeholders.
for i, slide in enumerate(prs.slides, 1):
    for shape in iter_text_shapes(slide):
        if not shape.text.strip() and '占位符' in shape.name:
            set_text(shape, f'{i:02d}', size=9, color=MUTED, bold=False, align=PP_ALIGN.RIGHT)

prs.save(OUT)
print(OUT)
